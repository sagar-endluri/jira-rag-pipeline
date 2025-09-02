import os
import requests
from requests.auth import HTTPBasicAuth
from openai import OpenAI
from dotenv import load_dotenv
import mimetypes
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import base64
from PIL import Image

class JiraAttachmentProcessor:
    def __init__(self):
        load_dotenv(".env")
        self.jira_url = os.getenv("jira_url")
        self.jira_api_token = os.getenv("jira_api_token")
        self.user_email = os.getenv("user_email")
        self.api_key = os.getenv("OPENAI_API_KEY")
        os.environ["OPENAI_API_KEY"] = self.api_key
        self.client = OpenAI(api_key=self.api_key)

        self.auth =  HTTPBasicAuth(self.user_email,self.jira_api_token )# "ATATT3xFfGF0hfORXaw3PFI1__5nKfH6fMc7F2oaMZm7Y7sFJ1K-Ip-vAe38Mpi52exFIf9qJQXTRenVv2k_gO02EYpji9pJysD4uA6Ucyca6lyZDJbdk2dL4c58Mf--Iq8LJTjeVTCNtvAKGiLc6H-3rNblZOt1LFPTqU7ERxEPGbOCBbbRpOU=4128C04E")

        self.headers = {"Accept": "*/*"}

    def download_attachment(self, url, filename, save_dir):
        os.makedirs(save_dir, exist_ok=True)
        file_path = os.path.join(save_dir, filename)
        text_path = os.path.join(save_dir, f"{filename}.txt")

        # try:
        response = requests.get(url, headers=self.headers, auth=self.auth, stream=True)
        if response.status_code == 200:
            with open(file_path, "wb") as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
            print(f"✅ Downloaded: {file_path}")
        else:
            raise Exception(f"Download failed with status code: {response.status_code}")

        extracted_text = self.extract_text_from_file(file_path, filename)

        if extracted_text:
            with open(text_path, "w", encoding="utf-8") as out:
                out.write(extracted_text)
            print(f"📝 Text extracted to: {text_path}")
            return extracted_text
        else:
            return "⚠️ No text extracted from file."

        # except Exception as e:
        #     print(f"❌ Error processing {filename}: {e}")
        #     return f"❌ Error: {e}"

    def extract_text_from_file(self, file_path, filename):
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return self.extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return self.extract_text_from_docx(file_path)
        elif ext == ".pptx":
            return self.extract_text_from_pptx(file_path)
        elif ext in [".png", ".jpg", ".jpeg"]:
            return self.extract_text_from_image(file_path)
        elif ext == ".txt":
            return open(file_path, "r", encoding="utf-8").read()
        else:
            print(f"📦 Unsupported file type for text extraction: {ext}")
            return None

    def extract_text_from_pdf(self, file_path):
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()

    def extract_text_from_docx(self, file_path):
        doc = DocxDocument(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    def extract_text_from_image(self, file_path):
        with open(file_path, "rb") as image_file:
            encoded_image = base64.b64encode(image_file.read()).decode()

        response = self.client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "user", "content": [
                    {"type": "text", "text": "Please describe the contents of this image."},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{encoded_image}"}}
                ]}
            ],
            max_tokens=1000
        )
        return response.choices[0].message.content
